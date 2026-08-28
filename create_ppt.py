"""
CTG-CPM Implementation Plan — Modern Dark-Themed PowerPoint Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── Color Palette (Dark Premium Theme) ──────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # #0D1117  near-black
BG_CARD      = RGBColor(0x16, 0x1B, 0x22)   # #161B22  card bg
ACCENT_BLUE  = RGBColor(0x58, 0xA6, 0xFF)   # #58A6FF  links / headings
ACCENT_GREEN = RGBColor(0x7E, 0xE7, 0x87)   # #7EE787  success
ACCENT_PURPLE= RGBColor(0xD2, 0xA8, 0xFF)   # #D2A8FF  gen-ai
ACCENT_ORANGE= RGBColor(0xFF, 0xA6, 0x57)   # #FFA657  warnings
ACCENT_RED   = RGBColor(0xFF, 0x7B, 0x72)   # #FF7B72  alerts
ACCENT_CYAN  = RGBColor(0x79, 0xC0, 0xFF)   # #79C0FF  secondary
TEXT_PRIMARY  = RGBColor(0xE6, 0xED, 0xF3)   # #E6EDF3  bright text
TEXT_SECONDARY= RGBColor(0x8B, 0x94, 0x9E)   # #8B949E  muted text
TEXT_DIM      = RGBColor(0x48, 0x4F, 0x58)   # #484F58  very muted
BORDER_COLOR  = RGBColor(0x30, 0x36, 0x3D)   # #30363D  borders
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRADIENT_START= RGBColor(0x0D, 0x11, 0x17)
GRADIENT_END  = RGBColor(0x14, 0x1D, 0x2B)

# ── Helpers ─────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD,
              border_color=BORDER_COLOR, border_width=Pt(1), corner_radius=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def set_para(tf, text, font_size=14, color=TEXT_PRIMARY, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Segoe UI", space_after=Pt(6),
             space_before=Pt(0), is_first=True):
    """Add a paragraph to a text frame."""
    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_circle_number(slide, left, top, size, number, color=ACCENT_BLUE):
    """Add a numbered circle badge."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return circle


def add_pill_badge(slide, left, top, text, bg_color=ACCENT_BLUE, text_color=BG_DARK, width=None):
    """Add a pill-shaped badge."""
    w = width or Inches(1.8)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_color
    pill.line.fill.background()
    pill.adjustments[0] = 0.5  # max rounding
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return pill


def add_slide_number(slide, number, total):
    """Add slide number indicator at bottom right."""
    tf = add_text_box(slide, Inches(11.0), Inches(7.1), Inches(2.0), Inches(0.35))
    set_para(tf, f"{number:02d} / {total:02d}", font_size=9, color=TEXT_DIM,
             alignment=PP_ALIGN.RIGHT, space_after=Pt(0))


def add_footer_bar(slide, slide_num, total):
    """Add a subtle footer bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.25), Inches(13.333), Pt(2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BORDER_COLOR
    bar.line.fill.background()
    # Slide number
    add_slide_number(slide, slide_num, total)
    # Brand
    tf = add_text_box(slide, Inches(0.6), Inches(7.1), Inches(3.0), Inches(0.35))
    set_para(tf, "CTG-CPM  //  IMPLEMENTATION PLAN", font_size=8, color=TEXT_DIM,
             alignment=PP_ALIGN.LEFT, space_after=Pt(0))


# ── Presentation Setup ─────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

TOTAL_SLIDES = 15

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Large accent rectangle on left
accent_block = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5)
)
accent_block.fill.solid()
accent_block.fill.fore_color.rgb = ACCENT_BLUE
accent_block.line.fill.background()

# Version badge
add_pill_badge(slide, Inches(1.0), Inches(1.0), "CTG-CPM  //  v1.0",
               bg_color=ACCENT_BLUE, text_color=BG_DARK, width=Inches(2.0))

# Patent badge
add_pill_badge(slide, Inches(3.2), Inches(1.0), "PATENT PENDING",
               bg_color=ACCENT_PURPLE, text_color=BG_DARK, width=Inches(1.6))

# GenAI badge
add_pill_badge(slide, Inches(5.0), Inches(1.0), "GenAI + AGENTIC AI + GAME THEORY",
               bg_color=ACCENT_GREEN, text_color=BG_DARK, width=Inches(3.2))

# Main title
tf = add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(2.0))
set_para(tf, "Predictive Maintenance for", font_size=42, color=TEXT_SECONDARY,
         bold=False, font_name="Segoe UI Light")
set_para(tf, "Counterfactual Telemetry", font_size=52, color=WHITE, bold=True,
         font_name="Segoe UI", is_first=False, space_before=Pt(4))

# Accent line
add_accent_line(slide, Inches(1.0), Inches(4.4), Inches(3.0), ACCENT_BLUE)

# Subtitle
tf = add_text_box(slide, Inches(1.0), Inches(4.7), Inches(10.0), Inches(1.2))
set_para(tf, "Recommendation-Driven Predictive Maintenance using Generative & Agentic AI",
         font_size=20, color=TEXT_SECONDARY, font_name="Segoe UI")
set_para(tf, "with Game-Theoretic Multi-Agent Optimization",
         font_size=20, color=ACCENT_PURPLE, font_name="Segoe UI", is_first=False)

# Bottom info
tf = add_text_box(slide, Inches(1.0), Inches(6.2), Inches(8.0), Inches(0.8))
set_para(tf, "IMPLEMENTATION PLAN  \u2022  AUGUST 2026  \u2022  CONFIDENTIAL",
         font_size=11, color=TEXT_DIM, font_name="Segoe UI")

# Decorative circles (top right)
for i, (x, y, sz, clr) in enumerate([
    (11.0, 0.5, 0.6, ACCENT_BLUE), (11.8, 1.0, 0.4, ACCENT_PURPLE),
    (11.5, 1.8, 0.3, ACCENT_GREEN), (12.2, 0.3, 0.35, ACCENT_ORANGE)
]):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                Inches(sz), Inches(sz))
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()
    # Set transparency via XML manipulation
    try:
        spPr = c._element.find(qn('a:solidFill'))
        if spPr is not None:
            srgb_elem = spPr.find(qn('a:srgbClr'))
            if srgb_elem is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb_elem, qn('a:alpha'))
                alpha.set('val', '30000')
    except Exception:
        pass  # Skip transparency if it fails


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABSTRACT
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 2, TOTAL_SLIDES)

# Section label
add_pill_badge(slide, Inches(0.8), Inches(0.5), "ABSTRACT", bg_color=ACCENT_BLUE, width=Inches(1.4))

# Title
tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.8))
set_para(tf, "Project Overview", font_size=36, color=WHITE, bold=True)

add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.0), ACCENT_BLUE)

# Abstract text card
card = add_shape(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.5),
                 fill_color=BG_CARD, border_color=BORDER_COLOR)

tf = add_text_box(slide, Inches(1.2), Inches(2.3), Inches(10.7), Inches(4.2))
abstract_text = (
    "Current predictive maintenance (PdM) systems in telecom and cloud infrastructure "
    "operate as opaque \u201Cblack boxes\u201D \u2014 they can predict that a failure will occur "
    "but cannot explain why or autonomously prescribe what to do about it."
)
set_para(tf, abstract_text, font_size=13, color=TEXT_SECONDARY, space_after=Pt(12))

abstract_text2 = (
    "CTG-CPM introduces a system that leverages Generative AI "
    "(Time-Series Diffusion Models) to project hypothetical future network telemetry, "
    "Graph Neural Networks (GNNs) to enforce topology-aware constraints, and Multi-Agent "
    "Agentic AI to produce remediation recommendations. "
    "Decision compute is sub-millisecond; commands are not auto-deployed in the prototype."
)
set_para(tf, abstract_text2, font_size=13, color=TEXT_PRIMARY, is_first=False, space_after=Pt(12))

abstract_text3 = (
    "We further integrate Algorithmic Game Theory concepts \u2014 VCG Auction Mechanisms, "
    "Shapley Value decomposition, and Nash Equilibrium coordination \u2014 to provably "
    "optimize the multi-agent decision-making pipeline."
)
set_para(tf, abstract_text3, font_size=13, color=ACCENT_PURPLE, is_first=False, space_after=Pt(12))

# Key metric badges at bottom
metrics = [
    ("FAST DECISION COMPUTE", ACCENT_GREEN),
    ("POTENTIAL OPEX SAVINGS", ACCENT_BLUE),
    ("RECOMMENDATION-BASED FIXES", ACCENT_PURPLE),
    ("3 PATENT CLAIMS", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(metrics):
    add_pill_badge(slide, Inches(0.8 + i * 2.9), Inches(6.7), label,
                   bg_color=color, text_color=BG_DARK, width=Inches(2.6))


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 3, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "THE PROBLEM", bg_color=ACCENT_RED, width=Inches(1.6))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.8))
set_para(tf, "Current Predictive Maintenance is a Black Box", font_size=34, color=WHITE, bold=True)

add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.5), ACCENT_RED)

# Three problem cards
problems = [
    ("\u26A0  The Alert", "Today\u2019s AI can predict that a network node will fail "
     "in X hours, but cannot explain WHY or WHAT to do about it.",
     ACCENT_ORANGE, "OPACITY"),
    ("\u26A1  The Risk", "NOC engineers must manually guess the best remediation "
     "and deploy it live, risking further cascading outages.",
     ACCENT_RED, "MANUAL FIX"),
    ("\U0001F4B0  The Cost", "Companies perform costly \u201Cblind\u201D hardware replacements "
     "(truck rolls) because they cannot safely test configurations on live traffic.",
     ACCENT_PURPLE, "OPEX DRAIN"),
]

for i, (title, desc, color, badge_text) in enumerate(problems):
    x = Inches(0.8 + i * 4.0)
    card = add_shape(slide, x, Inches(2.2), Inches(3.6), Inches(3.5),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Badge
    add_pill_badge(slide, x + Inches(0.3), Inches(2.45), badge_text,
                   bg_color=color, text_color=BG_DARK, width=Inches(1.5))
    
    # Number
    add_circle_number(slide, x + Inches(2.8), Inches(2.4), Inches(0.45),
                      str(i + 1), color)
    
    # Title
    tf = add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(3.0), Inches(0.5))
    set_para(tf, title, font_size=16, color=WHITE, bold=True)
    
    # Description
    tf = add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.0), Inches(1.8))
    set_para(tf, desc, font_size=12, color=TEXT_SECONDARY, space_after=Pt(0))

# Bottom comparison bar
bar = add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
                fill_color=RGBColor(0x1C, 0x10, 0x10), border_color=ACCENT_RED, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10.7), Inches(0.6))
set_para(tf, "CURRENT STATE:   AI Predicts Failure  \u2192  Human Writes Script  \u2192  "
         "MANUAL  \u2022  SLOW  \u2022  RISKY", font_size=12, color=ACCENT_RED, bold=True,
         alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 4, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "OUR SOLUTION", bg_color=ACCENT_GREEN, width=Inches(1.6))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.8))
set_para(tf, "Introducing CTG-CPM", font_size=36, color=WHITE, bold=True)

add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.0), ACCENT_GREEN)

tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(10.0), Inches(0.5))
set_para(tf, "A system that generates counterfactual futures, evaluates candidate "
         "remediations, and produces a RECOMMENDED command (not auto-deployed in the prototype).",
         font_size=14, color=TEXT_SECONDARY)

# Three solution steps
steps = [
    ("STEP 01", "GenAI for \u201CWhat-If\u201D Data",
     "We use Generative AI (Time-Series Diffusion) to synthesize \u201CCounterfactual "
     "Telemetry\u201D \u2014 simulated future network data based on hypothetical changes.",
     ACCENT_BLUE),
    ("STEP 02", "Agentic Causal Loop",
     "Multi-Agent AI analyzes this synthetic data to pinpoint exact root causes "
     "and test fixes. VCG auctions allocate tasks optimally.",
     ACCENT_PURPLE),
    ("STEP 03", "Remediation Recommendation",
     "The system selects a candidate via Nash Equilibrium coordination and "
     "produces a recommended command. In the prototype it is NOT auto-deployed.",
     ACCENT_GREEN),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 4.0)
    card = add_shape(slide, x, Inches(2.7), Inches(3.6), Inches(3.6),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Step badge
    add_pill_badge(slide, x + Inches(0.3), Inches(2.95), step,
                   bg_color=color, text_color=BG_DARK, width=Inches(1.2))
    
    # Circle number
    add_circle_number(slide, x + Inches(2.8), Inches(2.9), Inches(0.45),
                      str(i + 1), color)
    
    # Title
    tf = add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.0), Inches(0.5))
    set_para(tf, title, font_size=18, color=WHITE, bold=True)
    
    # Description
    tf = add_text_box(slide, x + Inches(0.3), Inches(4.1), Inches(3.0), Inches(1.8))
    set_para(tf, desc, font_size=12, color=TEXT_SECONDARY)

# Bottom bar
bar = add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                fill_color=RGBColor(0x0C, 0x1C, 0x0C), border_color=ACCENT_GREEN, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.2), Inches(6.55), Inches(10.7), Inches(0.5))
set_para(tf, "OUR APPROACH:   AI Detects  \u2192  AI Generates Counterfactual Futures  \u2192  "
         "AI Recommends a Remediation   \u2022  DecisioN Compute ms  \u2022  Not Auto-Deployed",
         font_size=12, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (4-STAGE PIPELINE)
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 5, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "HOW IT WORKS", bg_color=ACCENT_CYAN, width=Inches(1.6))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.8))
set_para(tf, "The 4-Stage Patentable Pipeline", font_size=34, color=WHITE, bold=True)

add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.5), ACCENT_CYAN)

tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(10.0), Inches(0.4))
set_para(tf, "From live telemetry ingestion to remediation recommendation \u2014 a closed pipeline that runs in milliseconds of decision compute (not deployment).",
         font_size=13, color=TEXT_SECONDARY)

stages = [
    ("01", "DETECT", "live", "GNN maps live topology\nand detects anomalies\n(thermal, optical, traffic)", ACCENT_BLUE),
    ("02", "GENERATE", "ms", "Diffusion model creates\nN counterfactual streams\nconditioned on K\ninterventions", ACCENT_PURPLE),
    ("03", "SELECT", "ms", "Agent 1: Root cause via\nShapley\nAgent 2: SPE + VCG\nAgent 3: pick via projection", ACCENT_GREEN),
    ("04", "RECOMMEND", "only", "Generates recommended\ncommand; NOT auto-deployed\nin the prototype", ACCENT_ORANGE),
]

for i, (num, name, time_est, desc, color) in enumerate(stages):
    x = Inches(0.8 + i * 3.1)
    
    # Card
    card = add_shape(slide, x, Inches(2.6), Inches(2.7), Inches(3.8),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Stage number circle
    add_circle_number(slide, x + Inches(0.3), Inches(2.85), Inches(0.5), num, color)
    
    # Stage name
    tf = add_text_box(slide, x + Inches(0.95), Inches(2.85), Inches(1.5), Inches(0.4))
    set_para(tf, name, font_size=18, color=color, bold=True)
    
    # Time badge
    add_pill_badge(slide, x + Inches(0.95), Inches(3.3), time_est,
                   bg_color=color, text_color=BG_DARK, width=Inches(0.8))
    
    # Description
    tf = add_text_box(slide, x + Inches(0.3), Inches(3.8), Inches(2.2), Inches(2.2))
    set_para(tf, desc, font_size=11, color=TEXT_SECONDARY)
    
    # Arrow between stages
    if i < 3:
        arrow_x = x + Inches(2.85)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(4.3), Inches(0.3), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_DIM
        arrow.line.fill.background()

# Total latency bar
bar = add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.45),
                fill_color=BG_CARD, border_color=ACCENT_CYAN, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.0), Inches(6.62), Inches(11.1), Inches(0.4))
set_para(tf, "IN-PROCESS DECISION COMPUTE:  milliseconds   "
         "\u2502  Excludes LLM round-trip and real device deployment; end-to-end MTTR not claimed",
         font_size=11, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — ALGORITHMS (Part 1: Core ML)
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 6, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ALGORITHMS", bg_color=ACCENT_PURPLE, width=Inches(1.4))
add_pill_badge(slide, Inches(2.4), Inches(0.5), "PART 1: CORE ML STACK", bg_color=BORDER_COLOR, text_color=TEXT_PRIMARY, width=Inches(2.4))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.8))
set_para(tf, "Core Algorithm Stack", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.0), ACCENT_PURPLE)

# Algorithm 1
card = add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.2),
                 fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(2))
add_circle_number(slide, Inches(1.1), Inches(2.3), Inches(0.45), "1", ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.7), Inches(2.3), Inches(4.3), Inches(0.4))
set_para(tf, "Time-Series Diffusion Model", font_size=16, color=ACCENT_BLUE, bold=True)
tf = add_text_box(slide, Inches(1.1), Inches(2.8), Inches(4.9), Inches(0.3))
set_para(tf, "Counterfactual Telemetry Generator", font_size=11, color=TEXT_DIM, bold=True)
tf = add_text_box(slide, Inches(1.1), Inches(3.1), Inches(4.9), Inches(1.0))
set_para(tf, "\u2022 Denoising score matching on multivariate time-series\n"
         "\u2022 Conditioned on intervention variables\n"
         "\u2022 Protocol-constrained generation (valid OSNR, IP ranges)\n"
         "\u2022 Based on Diffusion-TS architecture (ICLR 2024)",
         font_size=10, color=TEXT_SECONDARY)

# Algorithm 2
card = add_shape(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.2),
                 fill_color=BG_CARD, border_color=ACCENT_GREEN, border_width=Pt(2))
add_circle_number(slide, Inches(7.1), Inches(2.3), Inches(0.45), "2", ACCENT_GREEN)
tf = add_text_box(slide, Inches(7.7), Inches(2.3), Inches(4.3), Inches(0.4))
set_para(tf, "Graph Neural Network (GraphSAGE)", font_size=16, color=ACCENT_GREEN, bold=True)
tf = add_text_box(slide, Inches(7.1), Inches(2.8), Inches(4.9), Inches(0.3))
set_para(tf, "Topology-Aware Anomaly Detection", font_size=11, color=TEXT_DIM, bold=True)
tf = add_text_box(slide, Inches(7.1), Inches(3.1), Inches(4.9), Inches(1.0))
set_para(tf, "\u2022 GraphSAGE with temporal attention layers\n"
         "\u2022 Spatial (inter-device) + Temporal (time-series) learning\n"
         "\u2022 Cascade failure propagation modeling\n"
         "\u2022 Dynamic graph updates every 60 seconds",
         font_size=10, color=TEXT_SECONDARY)

# Algorithm 3
card = add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.2),
                 fill_color=BG_CARD, border_color=ACCENT_ORANGE, border_width=Pt(2))
add_circle_number(slide, Inches(1.1), Inches(4.8), Inches(0.45), "3", ACCENT_ORANGE)
tf = add_text_box(slide, Inches(1.7), Inches(4.8), Inches(5.0), Inches(0.4))
set_para(tf, "Multi-Agent Causal Loop (Agentic AI Orchestration)", font_size=16, color=ACCENT_ORANGE, bold=True)

# Three agents
agents = [
    ("Agent 1 \u2014 Diagnostician", "Granger Causality + Do-Calculus\nfor root cause identification", ACCENT_BLUE),
    ("Agent 2 \u2014 Projector", "Runs counterfactual diffusion\nto project 'what-if' outcomes", ACCENT_GREEN),
    ("Agent 3 \u2014 Advisor", "Generates recommended\nNETCONF/YANG config draft", ACCENT_ORANGE),
]
for i, (name, desc, clr) in enumerate(agents):
    x = Inches(1.1 + i * 3.7)
    add_pill_badge(slide, x, Inches(5.4), name, bg_color=clr, text_color=BG_DARK, width=Inches(2.8))
    tf = add_text_box(slide, x, Inches(5.8), Inches(3.2), Inches(0.8))
    set_para(tf, desc, font_size=10, color=TEXT_SECONDARY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — ALGORITHMS (Part 2: Game Theory)
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 7, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ALGORITHMS", bg_color=ACCENT_PURPLE, width=Inches(1.4))
add_pill_badge(slide, Inches(2.4), Inches(0.5), "PART 2: GAME THEORY", bg_color=ACCENT_ORANGE, text_color=BG_DARK, width=Inches(2.4))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.8))
set_para(tf, "Algorithmic Game Theory Integration", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.75), Inches(2.5), ACCENT_ORANGE)

tf = add_text_box(slide, Inches(0.8), Inches(1.95), Inches(10.0), Inches(0.35))
set_para(tf, "Provably optimal multi-agent coordination \u2014 not heuristic, not approximate, mathematically guaranteed.",
         font_size=12, color=TEXT_SECONDARY)

# Four game theory modules
gt_modules = [
    ("VCG Auction", "Agent Task Allocation",
     "Agents bid on remediation tasks based on true capability. "
     "VCG guarantees DSIC \u2014 truthful bidding is the dominant strategy, "
     "ensuring globally optimal assignment.",
     ACCENT_BLUE, "DSIC GUARANTEED"),
    ("Shapley Value", "Root-Cause Attribution",
     "Decomposes anomaly contribution across telemetry features "
     "with axiomatic fairness: Efficiency, Symmetry, Linearity, "
     "and Null-Player properties.",
     ACCENT_GREEN, "AXIOMATIC FAIR"),
    ("Nash Equilibrium", "Conflict Resolution",
     "Models remediation as a non-cooperative game. "
     "Fictitious Play converges to NE \u2014 no agent can unilaterally "
     "improve by changing its action.",
     ACCENT_PURPLE, "STABLE STRATEGY"),
    ("Groves Mechanism", "Truthful Telemetry",
     "In multi-vendor networks, incentivizes domains to "
     "report truthful telemetry. Payment depends on others\u2019 "
     "welfare \u2192 truth is dominant strategy.",
     ACCENT_ORANGE, "TRUTH-DOMINANT"),
]

for i, (title, subtitle, desc, color, badge) in enumerate(gt_modules):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(2.5 + row * 2.3)
    
    card = add_shape(slide, x, y, Inches(5.6), Inches(2.0),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    add_circle_number(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.4),
                      str(i + 1), color)
    
    tf = add_text_box(slide, x + Inches(0.8), y + Inches(0.15), Inches(3.0), Inches(0.35))
    set_para(tf, title, font_size=16, color=color, bold=True)
    
    add_pill_badge(slide, x + Inches(3.8), y + Inches(0.2), badge,
                   bg_color=color, text_color=BG_DARK, width=Inches(1.6))
    
    tf = add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.0), Inches(0.25))
    set_para(tf, subtitle, font_size=10, color=TEXT_DIM, bold=True)
    
    tf = add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(5.0), Inches(1.0))
    set_para(tf, desc, font_size=11, color=TEXT_SECONDARY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — ADVANTAGES
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 8, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ADVANTAGES", bg_color=ACCENT_GREEN, width=Inches(1.5))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "Key Advantages of CTG-CPM", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.65), Inches(2.0), ACCENT_GREEN)

advantages = [
    ("Fast Decision Compute", "In-process compute is sub-millisecond (excludes LLM & deployment)", ACCENT_GREEN),
    ("Projection-based Selection", "Remediation chosen from projected scenarios; not auto-deployed", ACCENT_BLUE),
    ("Potential OPEX Benefit", "May avoid truck-rolls (not yet measured)", ACCENT_ORANGE),
    ("Explainable Attribution", "Shapley attribution for transparency", ACCENT_PURPLE),
    ("Data Scarcity Fix", "GenAI creates unlimited synthetic data", ACCENT_CYAN),
    ("Topology-Aware", "GNN respects physical network constraints", ACCENT_GREEN),
    ("Optimal Coordination", "Game theory guarantees conflict-free ops", ACCENT_ORANGE),
    ("Vendor-Agnostic", "NETCONF/YANG works across all vendors", ACCENT_BLUE),
    ("Scalable Agents", "Add more agents without architecture changes", ACCENT_PURPLE),
    ("Patent Protected", "3 novel patent claims for IP defense", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(advantages):
    row = i // 5
    col = i % 5
    x = Inches(0.8 + col * 2.4)
    y = Inches(2.0 + row * 2.5)
    
    card = add_shape(slide, x, y, Inches(2.1), Inches(2.2),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    
    add_circle_number(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.38),
                      str(i + 1), color)
    
    tf = add_text_box(slide, x + Inches(0.15), y + Inches(0.65), Inches(1.8), Inches(0.45))
    set_para(tf, title, font_size=13, color=color, bold=True)
    
    tf = add_text_box(slide, x + Inches(0.15), y + Inches(1.15), Inches(1.8), Inches(0.9))
    set_para(tf, desc, font_size=10, color=TEXT_SECONDARY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — DISADVANTAGES
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 9, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "LIMITATIONS", bg_color=ACCENT_RED, width=Inches(1.5))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "Challenges & Mitigation Strategies", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT_RED)

disadvantages = [
    ("Computational Cost", "Diffusion models are iteratively expensive",
     "Distilled diffusion (fewer steps) + edge-optimized architectures"),
    ("Cold-Start Problem", "Requires initial telemetry for GNN training",
     "Transfer learning from simulated environments (Simu5G)"),
    ("Sim-Reality Gap", "Digital twin may not perfectly replicate live behavior",
     "Domain randomization + periodic model recalibration"),
    ("Regulatory Risk", "Applying recommended remediations on critical infrastructure",
     "Human-in-the-loop approval + full audit trails"),
    ("Multi-Vendor Complexity", "NETCONF/YANG varies across vendors",
     "Vendor abstraction layer with per-vendor adaptors"),
    ("Game Theory Overhead", "VCG computation is NP-hard in general",
     "Bounded agent/task spaces + polynomial approximations"),
    ("Adversarial Robustness", "Counterfactuals could be manipulated",
     "Adversarial training + anomaly detection on synthetic quality"),
    ("Trust & Adoption", "Engineers may resist fully autonomous systems",
     "Graduated autonomy: Advisory \u2192 Semi-Auto \u2192 Fully Auto"),
]

for i, (title, issue, mitigation) in enumerate(disadvantages):
    row = i // 4
    col = i % 4
    x = Inches(0.8 + col * 3.05)
    y = Inches(2.0 + row * 2.55)
    
    card = add_shape(slide, x, y, Inches(2.75), Inches(2.3),
                     fill_color=BG_CARD, border_color=ACCENT_RED, border_width=Pt(1))
    
    add_circle_number(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.35),
                      str(i + 1), ACCENT_RED)
    
    tf = add_text_box(slide, x + Inches(0.6), y + Inches(0.12), Inches(2.0), Inches(0.35))
    set_para(tf, title, font_size=12, color=ACCENT_RED, bold=True)
    
    tf = add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.45), Inches(0.6))
    set_para(tf, issue, font_size=9, color=TEXT_SECONDARY)
    
    # Mitigation
    mit_bar = add_shape(slide, x + Inches(0.1), y + Inches(1.2), Inches(2.55), Inches(0.95),
                        fill_color=RGBColor(0x0F, 0x1A, 0x0F), border_color=ACCENT_GREEN, border_width=Pt(1))
    tf = add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(0.2))
    set_para(tf, "\u2713 MITIGATION", font_size=8, color=ACCENT_GREEN, bold=True)
    tf = add_text_box(slide, x + Inches(0.2), y + Inches(1.45), Inches(2.35), Inches(0.65))
    set_para(tf, mitigation, font_size=9, color=ACCENT_GREEN)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — LITERATURE REVIEW (Table 1 of 2)
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 10, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "LITERATURE REVIEW", bg_color=ACCENT_CYAN, width=Inches(1.8))
add_pill_badge(slide, Inches(2.8), Inches(0.5), "PAGE 1 OF 2", bg_color=BORDER_COLOR, text_color=TEXT_PRIMARY, width=Inches(1.3))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "Related Work & Differentiation", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(2.0), ACCENT_CYAN)

# Table header
header_y = Inches(1.85)
headers = [
    (Inches(0.4), Inches(0.3), "#"),
    (Inches(0.75), Inches(2.8), "Paper / Authors"),
    (Inches(3.6), Inches(1.0), "Year"),
    (Inches(4.65), Inches(1.8), "Venue"),
    (Inches(6.5), Inches(2.2), "Algorithm"),
    (Inches(8.75), Inches(3.7), "CTG-CPM Differentiator"),
]

header_bar = add_shape(slide, Inches(0.3), header_y, Inches(12.5), Inches(0.35),
                       fill_color=ACCENT_BLUE, border_color=ACCENT_BLUE)
for x, w, text in headers:
    tf = add_text_box(slide, x, header_y, w, Inches(0.35))
    set_para(tf, text, font_size=9, color=BG_DARK, bold=True, alignment=PP_ALIGN.LEFT)

# Table rows (papers 1-8)
papers_1 = [
    ("1", "Diffusion-TS\nZeng, Chen, Zhang, Xu", "2024", "ICLR", "Denoising Diffusion\n+ Transformer", "We add GNN topology\nconstraints + intervention\nconditioning"),
    ("2", "TSDiff\nKollovieh, Ansari et al.", "2023", "NeurIPS", "Unconditional Diffusion\n+ Self-guidance", "CTG-CPM uses conditional\ngeneration tied to\nremediation actions"),
    ("3", "NetDiffusion\nJiang, Liu, Feamster et al.", "2024", "ACM\nPOMACS", "Protocol-aware\nStable Diffusion", "We generate telemetry-\nlevel counterfactuals,\nnot packet-level traffic"),
    ("4", "Simba (5G RCA)\nHasan, Boeira et al.", "2024", "arXiv", "GNN +\nTransformer", "Simba detects only;\nCTG-CPM adds\nremediation recommendations"),
    ("5", "Time-Series Diffusion\nSurvey — Yang et al.", "2024", "ACM\nComp.\nSurveys", "Survey of DDPM,\nScore-based models", "Survey only; CTG-CPM\nadds counterfactual\nremediation selection"),
    ("6", "Autohma-LLM\nIEEE TCCN authors", "2025", "IEEE\nTCCN", "LLM hybrid\nmulti-agent", "We use VCG auctions\nfor provable optimality\nvs. heuristic LLM prompts"),
    ("7", "GNN+SHAP for 5G SDN\nSciOpen / Tsinghua", "2024", "Big Data\nMining &\nAnalytics", "GNN + Multi-Head\nAttention + SHAP", "Shapley used in-loop\nfor real-time RCA,\nnot post-hoc explanation"),
    ("8", "VCG+MARL for 6G IoV\narXiv cs.GT / cs.NI", "2025", "arXiv", "VCG + MARL\nfor resource slicing", "VCG applied to\nremediation allocation,\nnot traffic management"),
]

for i, (num, paper, year, venue, algo, diff) in enumerate(papers_1):
    y = Inches(2.25 + i * 0.62)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    
    row_bar = add_shape(slide, Inches(0.3), y, Inches(12.5), Inches(0.58),
                        fill_color=bg, border_color=BORDER_COLOR, border_width=Pt(0.5))
    
    cols = [
        (Inches(0.4), Inches(0.3), num, 9, TEXT_DIM),
        (Inches(0.75), Inches(2.8), paper, 8, TEXT_PRIMARY),
        (Inches(3.6), Inches(1.0), year, 9, ACCENT_BLUE),
        (Inches(4.65), Inches(1.8), venue, 8, TEXT_SECONDARY),
        (Inches(6.5), Inches(2.2), algo, 8, ACCENT_PURPLE),
        (Inches(8.75), Inches(3.7), diff, 8, ACCENT_GREEN),
    ]
    for cx, cw, ctext, csize, ccolor in cols:
        tf = add_text_box(slide, cx, y + Inches(0.02), cw, Inches(0.55))
        set_para(tf, ctext, font_size=csize, color=ccolor, space_after=Pt(0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — LITERATURE REVIEW (Table 2 of 2)
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 11, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "LITERATURE REVIEW", bg_color=ACCENT_CYAN, width=Inches(1.8))
add_pill_badge(slide, Inches(2.8), Inches(0.5), "PAGE 2 OF 2", bg_color=BORDER_COLOR, text_color=TEXT_PRIMARY, width=Inches(1.3))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "Related Work & Differentiation (cont.)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(2.5), ACCENT_CYAN)

# Table header
header_y = Inches(1.85)
header_bar = add_shape(slide, Inches(0.3), header_y, Inches(12.5), Inches(0.35),
                       fill_color=ACCENT_BLUE, border_color=ACCENT_BLUE)
for x, w, text in headers:
    tf = add_text_box(slide, x, header_y, w, Inches(0.35))
    set_para(tf, text, font_size=9, color=BG_DARK, bold=True, alignment=PP_ALIGN.LEFT)

papers_2 = [
    ("9", "Latency Anomaly in 5G\nMDPI Sensors authors", "2024", "MDPI\nSensors", "ConvAE +\nLSTM Transfer", "Limited to latency;\nCTG-CPM handles\nmultivariate telemetry"),
    ("10", "Robust Adaptive Mechs.\narXiv cs.GT authors", "2025", "arXiv", "Online Learning\n+ VCG-style", "Theoretical only;\nCTG-CPM applies in\nconcrete network system"),
    ("11", "GAPPO Task Allocation\nMDPI Applied Sciences", "2024", "MDPI\nApplied\nSciences", "Genetic Algo\n+ PPO (MARL)", "We use formal game\ntheory guarantees, not\nevolutionary heuristics"),
    ("12", "MARL Network Healing\nIEEE Trans. Veh. Tech.", "2024", "IEEE\nTVT", "MAPPO,\nMADDPG", "CTG-CPM agents use\nsynthetic counterfactual\ndata, not just live obs."),
    ("13", "Generative DT for PdM\nMDPI Machines", "2024", "MDPI\nMachines", "GAN-based\nDigital Twin", "Diffusion > GAN quality;\nplus agentic recommendation layer not in scope"),
    ("14", "Causal AI & Digital Twins\nIEEE/ACM ASE + RESS", "2025", "IEEE/ACM\nASE", "Causal DAGs,\nDo-Calculus", "Theoretical framework;\nCTG-CPM implements a\ncausal attribution loop"),
    ("15", "GenAI Edge PdM\nPothireddy N.K.R.", "2025", "IEEE\nTransactions", "Lightweight GenAI\nfor IoT edge", "Edge-focused only;\nCTG-CPM spans full\nnetwork + game theory"),
]

for i, (num, paper, year, venue, algo, diff) in enumerate(papers_2):
    y = Inches(2.25 + i * 0.65)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    
    row_bar = add_shape(slide, Inches(0.3), y, Inches(12.5), Inches(0.6),
                        fill_color=bg, border_color=BORDER_COLOR, border_width=Pt(0.5))
    
    cols = [
        (Inches(0.4), Inches(0.3), num, 9, TEXT_DIM),
        (Inches(0.75), Inches(2.8), paper, 8, TEXT_PRIMARY),
        (Inches(3.6), Inches(1.0), year, 9, ACCENT_BLUE),
        (Inches(4.65), Inches(1.8), venue, 8, TEXT_SECONDARY),
        (Inches(6.5), Inches(2.2), algo, 8, ACCENT_PURPLE),
        (Inches(8.75), Inches(3.7), diff, 8, ACCENT_GREEN),
    ]
    for cx, cw, ctext, csize, ccolor in cols:
        tf = add_text_box(slide, cx, y + Inches(0.02), cw, Inches(0.55))
        set_para(tf, ctext, font_size=csize, color=ccolor, space_after=Pt(0))

# Summary badge at bottom
bar = add_shape(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.35),
                fill_color=BG_CARD, border_color=ACCENT_CYAN, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.0), Inches(6.85), Inches(11.1), Inches(0.35))
set_para(tf, "15 PAPERS REVIEWED  \u2502  7 CONFERENCES/JOURNALS  \u2502  "
         "4 NOVEL DIFFERENTIATORS: Counterfactual GenAI + Agentic Loop + GNN Topology + Game Theory",
         font_size=9, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 12, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ARCHITECTURE", bg_color=ACCENT_PURPLE, width=Inches(1.6))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "5-Layer System Architecture", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(2.0), ACCENT_PURPLE)

# 5 Layer boxes (stacked horizontally)
layers = [
    ("LAYER 1", "DATA", "Kafka, gNMI\nInfluxDB, Neo4j\netcd", ACCENT_BLUE),
    ("LAYER 2", "TOPOLOGY", "GraphSAGE GNN\nAnomaly Detector\nCascade Predictor", ACCENT_GREEN),
    ("LAYER 3", "GenAI", "Time-Series\nDiffusion Model\nCounterfactual Gen.", ACCENT_PURPLE),
    ("LAYER 4", "AGENTIC", "3 AI Agents\nVCG Auction\nNash Equilibrium", ACCENT_ORANGE),
    ("LAYER 5", "EXECUTION", "NETCONF/YANG\nConfig Validator\nRollback Engine", ACCENT_RED),
]

for i, (layer_num, layer_name, tech, color) in enumerate(layers):
    x = Inches(0.5 + i * 2.5)
    
    card = add_shape(slide, x, Inches(2.0), Inches(2.2), Inches(4.2),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Layer number badge
    add_pill_badge(slide, x + Inches(0.15), Inches(2.2), layer_num,
                   bg_color=color, text_color=BG_DARK, width=Inches(1.0))
    
    # Layer name
    tf = add_text_box(slide, x + Inches(0.15), Inches(2.65), Inches(1.9), Inches(0.5))
    set_para(tf, layer_name, font_size=22, color=color, bold=True)
    
    # Separator line
    add_accent_line(slide, x + Inches(0.15), Inches(3.15), Inches(1.9), color)
    
    # Tech details
    tf = add_text_box(slide, x + Inches(0.15), Inches(3.35), Inches(1.9), Inches(2.5))
    for j, line in enumerate(tech.split("\n")):
        set_para(tf, "\u2022 " + line, font_size=10, color=TEXT_SECONDARY,
                 is_first=(j == 0), space_after=Pt(4))
    
    # Arrow between layers
    if i < 4:
        arrow_x = x + Inches(2.3)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(3.8), Inches(0.25), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_DIM
        arrow.line.fill.background()

# Bottom: Data flow
bar = add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.0), Inches(0.5),
                fill_color=BG_CARD, border_color=ACCENT_PURPLE, border_width=Pt(1))
tf = add_text_box(slide, Inches(0.7), Inches(6.42), Inches(11.6), Inches(0.45))
set_para(tf, "DATA FLOW:  Live Telemetry  \u2192  Graph Mapping  \u2192  "
         "Counterfactual Projection  \u2192  Agent Selection  \u2192  Remediation Recommendation (not deployed)   "
         "\u2502  5 LAYERS",
         font_size=10, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — USE CASE
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 13, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "USE CASE", bg_color=ACCENT_ORANGE, width=Inches(1.3))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "5G Backhaul Fiber Link Failure", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(2.0), ACCENT_ORANGE)

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(10.0), Inches(0.4))
set_para(tf, "A 5G optical transceiver shows micro-fluctuations in OSNR. Standard PdM flags it for hardware replacement in 2 days.",
         font_size=12, color=TEXT_SECONDARY)

# 4 steps
use_steps = [
    ("01", "DETECT", "live",
     "GNN detects OSNR anomaly.\nCascade Risk Predictor flags\n3 downstream nodes at risk.",
     ACCENT_BLUE),
    ("02", "GENERATE", "ms",
     "Diffusion generates 5 counter-\nfactual streams: shift load,\nincrease laser bias, reroute,\nreduce traffic, replace HW.",
     ACCENT_PURPLE),
    ("03", "SELECT", "ms",
     "VCG assigns Agent 1 to diagnose.\nShapley: Laser = 72% cause.\nAgent 2 picks the best\nprojection (not time-validated).",
     ACCENT_GREEN),
    ("04", "RECOMMEND", "only",
     "Agent 3 generates a NETCONF/\nYANG config change for laser\nbias as a RECOMMENDATION.\nNot auto-deployed.",
     ACCENT_ORANGE),
]

for i, (num, name, time_est, desc, color) in enumerate(use_steps):
    x = Inches(0.8 + i * 3.05)
    card = add_shape(slide, x, Inches(2.4), Inches(2.7), Inches(3.0),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    add_circle_number(slide, x + Inches(0.2), Inches(2.6), Inches(0.4), num, color)
    tf = add_text_box(slide, x + Inches(0.7), Inches(2.58), Inches(1.3), Inches(0.35))
    set_para(tf, name, font_size=15, color=color, bold=True)
    add_pill_badge(slide, x + Inches(1.8), Inches(2.6), time_est,
                   bg_color=color, text_color=BG_DARK, width=Inches(0.7))
    
    tf = add_text_box(slide, x + Inches(0.2), Inches(3.15), Inches(2.3), Inches(2.0))
    set_para(tf, desc, font_size=10, color=TEXT_SECONDARY)
    
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.8), Inches(3.7), Inches(0.3), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_DIM
        arrow.line.fill.background()

# Results bar
results = [
    ("ZERO DOWNTIME", ACCENT_GREEN),
    ("ZERO TRUCK ROLL", ACCENT_BLUE),
    ("6-MONTH LIFE EXTENSION", ACCENT_PURPLE),
    ("FULL AUDIT TRAIL", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(results):
    add_pill_badge(slide, Inches(0.8 + i * 3.0), Inches(5.7), label,
                   bg_color=color, text_color=BG_DARK, width=Inches(2.7))

# Context bar
bar = add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
                fill_color=RGBColor(0x0C, 0x1C, 0x0C), border_color=ACCENT_GREEN, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.0), Inches(6.25), Inches(11.1), Inches(0.5))
set_para(tf, "RESULT:  Remediation RECOMMENDATION generated. In the prototype the command is NOT "
         "auto-deployed; no service-interruption or equipment-life claims are made without live validation.",
         font_size=11, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — EXPECTED IMPACT / KPIs
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 14, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "IMPACT & KPIs", bg_color=ACCENT_GREEN, width=Inches(1.5))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "Transforming Network Operations", font_size=34, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(2.0), ACCENT_GREEN)

# 4 big KPI cards
kpis = [
    ("< few", "ms", "Decision Compute", "In-process compute only\n(excludes LLM & deployment)", ACCENT_GREEN),
    ("POTENTIAL", "OPEX", "Truck-Roll Savings", "Not yet measured;\nprojected benefit only", ACCENT_BLUE),
    ("RECOMMEND", "ONLY", "Safety Model", "Commands are recommended,\nnot auto-deployed", ACCENT_PURPLE),
    ("0", "CASCADING", "Outages Prevented", "Claim—not yet validated\non live infrastructure", ACCENT_ORANGE),
]

for i, (big_num, unit, title, desc, color) in enumerate(kpis):
    x = Inches(0.8 + i * 3.05)
    card = add_shape(slide, x, Inches(2.0), Inches(2.7), Inches(2.8),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Big number
    tf = add_text_box(slide, x + Inches(0.2), Inches(2.15), Inches(2.3), Inches(0.8))
    set_para(tf, big_num, font_size=44, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Unit
    tf = add_text_box(slide, x + Inches(0.2), Inches(2.85), Inches(2.3), Inches(0.3))
    set_para(tf, unit, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Separator
    add_accent_line(slide, x + Inches(0.3), Inches(3.2), Inches(2.1), color)
    
    # Title
    tf = add_text_box(slide, x + Inches(0.2), Inches(3.35), Inches(2.3), Inches(0.35))
    set_para(tf, title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Description
    tf = add_text_box(slide, x + Inches(0.2), Inches(3.75), Inches(2.3), Inches(0.8))
    set_para(tf, desc, font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Additional KPIs bar
bar = add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5),
                fill_color=BG_CARD, border_color=BORDER_COLOR)
tf = add_text_box(slide, Inches(1.2), Inches(5.2), Inches(3.0), Inches(0.3))
set_para(tf, "ADDITIONAL KPIs", font_size=10, color=TEXT_DIM, bold=True)

extra_kpis = [
    ("Counterfactual Quality", "TARGET FID < 50", "Target only; measured FID is currently far higher (poor on synthetic prior)"),
    ("Game-Theory Optimality", "VCG-maximizing", "VCG assignment maximizes social welfare of stated bids (exact on given bids)"),
    ("Decision Compute", "ms", "In-process compute only; excludes LLM & deployment"),
]
for i, (name, value, desc) in enumerate(extra_kpis):
    x = Inches(1.2 + i * 3.7)
    tf = add_text_box(slide, x, Inches(5.5), Inches(3.3), Inches(0.3))
    set_para(tf, name, font_size=11, color=TEXT_PRIMARY, bold=True)
    
    add_pill_badge(slide, x, Inches(5.85), value, bg_color=ACCENT_CYAN, text_color=BG_DARK, width=Inches(1.0))
    
    tf = add_text_box(slide, x + Inches(1.1), Inches(5.85), Inches(2.5), Inches(0.3))
    set_para(tf, desc, font_size=9, color=TEXT_SECONDARY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — ROADMAP & NEXT STEPS
# ════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 15, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ROADMAP", bg_color=ACCENT_BLUE, width=Inches(1.3))

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.0), Inches(0.7))
set_para(tf, "The Road to Autonomy", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.65), Inches(2.0), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(1.85), Inches(8.0), Inches(0.35))
set_para(tf, "6-MONTH PATH TO PILOT", font_size=12, color=TEXT_DIM, bold=True)

# 3 Phase cards
phases = [
    ("PHASE 01", "MONTHS 1\u20132", "FOUNDATIONS",
     "\u2022 Prior Art Search & PPA Filing\n"
     "\u2022 Complete Literature Review\n"
     "\u2022 Architecture Design Finalized\n"
     "\u2022 Data Collection Pipeline Setup\n"
     "\u2022 Team Assembly & Tooling",
     ACCENT_BLUE, "COMPLETED"),
    ("PHASE 02", "MONTHS 2\u20134", "PROTOTYPE",
     "\u2022 Build all 5 layers in sandbox\n"
     "\u2022 Diffusion model training (Simu5G)\n"
     "\u2022 GNN topology engine integration\n"
     "\u2022 Multi-Agent orchestration\n"
     "\u2022 Game Theory modules (VCG, Nash)",
     ACCENT_PURPLE, "CURRENT FOCUS"),
    ("PHASE 03", "MONTHS 4\u20136", "PRODUCTION",
     "\u2022 Live enterprise testbed pilot\n"
     "\u2022 Measure all KPIs in production\n"
     "\u2022 Iterate on performance\n"
     "\u2022 Regulatory compliance check\n"
     "\u2022 Prepare for scale deployment",
     ACCENT_GREEN, "UPCOMING"),
]

for i, (phase, timeline, subtitle, deliverables, color, status) in enumerate(phases):
    x = Inches(0.8 + i * 4.0)
    card = add_shape(slide, x, Inches(2.3), Inches(3.7), Inches(4.2),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Phase badge
    add_pill_badge(slide, x + Inches(0.25), Inches(2.5), phase,
                   bg_color=color, text_color=BG_DARK, width=Inches(1.2))
    
    # Status badge
    status_bg = ACCENT_GREEN if status == "COMPLETED" else (ACCENT_PURPLE if status == "CURRENT FOCUS" else BORDER_COLOR)
    add_pill_badge(slide, x + Inches(1.6), Inches(2.5), status,
                   bg_color=status_bg, text_color=BG_DARK if status != "UPCOMING" else TEXT_PRIMARY,
                   width=Inches(1.7))
    
    # Timeline
    tf = add_text_box(slide, x + Inches(0.25), Inches(3.0), Inches(3.2), Inches(0.3))
    set_para(tf, timeline, font_size=12, color=TEXT_DIM, bold=True)
    
    # Subtitle
    tf = add_text_box(slide, x + Inches(0.25), Inches(3.3), Inches(3.2), Inches(0.4))
    set_para(tf, subtitle, font_size=20, color=color, bold=True)
    
    # Separator
    add_accent_line(slide, x + Inches(0.25), Inches(3.75), Inches(3.2), color)
    
    # Deliverables
    tf = add_text_box(slide, x + Inches(0.25), Inches(3.9), Inches(3.2), Inches(2.3))
    set_para(tf, deliverables, font_size=11, color=TEXT_SECONDARY)
    
    # Arrow between phases
    if i < 2:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(3.8), Inches(4.2), Inches(0.25), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_DIM
        arrow.line.fill.background()

# Bottom tag
bar = add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
                fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
tf = add_text_box(slide, Inches(1.0), Inches(6.72), Inches(11.1), Inches(0.35))
set_para(tf, "TECHNOLOGY STACK:  Kafka \u2022 InfluxDB \u2022 Neo4j \u2022 PyTorch Geometric \u2022 "
         "Diffusion-TS \u2022 DoWhy \u2022 Nashpy \u2022 LangGraph \u2022 NETCONF \u2022 Grafana \u2022 MLflow",
         font_size=9, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════

output_path = os.environ.get("OUTPUT_PPT_PATH", r"d:\Predictive Maintenance Project 3\CTG-CPM_Implementation_Plan.pptx")
prs.save(output_path)
print(f"\n{'='*60}")
print(f"  SUCCESS: Presentation saved to:")
print(f"  {output_path}")
print(f"  {TOTAL_SLIDES} slides generated")
print(f"{'='*60}")
