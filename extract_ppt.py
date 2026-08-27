import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation(r'd:\Predictive Maintenance Project 3\CTG-CPM- Self-Healing Networks via Counterfactual Telemetry (1).pptx')

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}")
    print(f"{'='*60}")
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.encode('ascii', 'replace').decode('ascii')
            print(text)
        
        if shape.has_table:
            table = shape.table
            print("\n[TABLE]")
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.encode('ascii', 'replace').decode('ascii') for cell in row.cells]
                print(" | ".join(cells))
            print("[/TABLE]\n")
