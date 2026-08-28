"""
generate_review2_ppt.py
Generates a 15-Slide Presentation for Review 2: Design, Development & Progress
CTG-CPM: Self-Healing Networks via Counterfactual Telemetry & Game-Theoretic Multi-Agent Optimization
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Color Palette (Dark Theme) ──────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x11, 0x17)   # Near-black #0D1117
BG_CARD       = RGBColor(0x16, 0x1B, 0x22)   # Card bg #161B22
ACCENT_BLUE   = RGBColor(0x58, 0xA6, 0xFF)   # Links / headings #58A6FF
ACCENT_GREEN  = RGBColor(0x7E, 0xE7, 0x87)   # Success #7EE787
ACCENT_PURPLE = RGBColor(0xD2, 0xA8, 0xFF)   # GenAI #D2A8FF
ACCENT_ORANGE = RGBColor(0xFF, 0xA6, 0x57)   # Warnings #FFA657
ACCENT_RED    = RGBColor(0xFF, 0x7B, 0x72)   # Alerts #FF7B72
ACCENT_CYAN   = RGBColor(0x79, 0xC0, 0xFF)   # Secondary #79C0FF
TEXT_PRIMARY   = RGBColor(0xE6, 0xED, 0xF3)   # Bright text
TEXT_SECONDARY = RGBColor(0x8B, 0x94, 0x9E)   # Muted text
TEXT_DIM       = RGBColor(0x48, 0x4F, 0x58)   # Very muted
BORDER_COLOR   = RGBColor(0x30, 0x36, 0x3D)   # Borders
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=BG_CARD,
              border_color=BORDER_COLOR, border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def set_para(tf, text, font_size=14, color=TEXT_PRIMARY, bold=False,
              alignment=PP_ALIGN.LEFT, font_name="Segoe UI", space_after=Pt(4),
              space_before=Pt(0), is_first=True):
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_circle_number(slide, left, top, size, number, color=ACCENT_BLUE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return circle

def add_pill_badge(slide, left, top, text, bg_color=ACCENT_BLUE, text_color=BG_DARK, width=None):
    w = width or Inches(1.8)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_color
    pill.line.fill.background()
    pill.adjustments[0] = 0.5
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return pill

def add_footer_bar(slide, slide_num, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.25), Inches(13.333), Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BORDER_COLOR
    bar.line.fill.background()
    
    tf = add_text_box(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.35))
    set_para(tf, f"{slide_num:02d} / {total:02d}", font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)
    
    tf = add_text_box(slide, Inches(0.6), Inches(7.05), Inches(5.0), Inches(0.35))
    set_para(tf, "CTG-CPM  //  REVIEW 2: DESIGN, DEVELOPMENT & PROGRESS", font_size=8, color=TEXT_DIM, alignment=PP_ALIGN.LEFT)

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL_SLIDES = 15

# ====================================================================
# SLIDE 1: TITLE SLIDE
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_pill_badge(slide, Inches(0.8), Inches(0.8), "REVIEW 2 PRESENTATION", bg_color=ACCENT_BLUE, text_color=BG_DARK, width=Inches(2.2))
add_pill_badge(slide, Inches(3.2), Inches(0.8), "DESIGN, DEVELOPMENT & PROGRESS", bg_color=ACCENT_GREEN, text_color=BG_DARK, width=Inches(3.2))

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.2))
set_para(tf, "CTG-CPM: Self-Healing Networks via Counterfactual Telemetry", font_size=36, color=WHITE, bold=True)
set_para(tf, "Predictive Maintenance Recommendation using Generative AI, Game Theory & Multi-Agent Orchestration",
         font_size=18, color=ACCENT_CYAN, is_first=False, space_before=Pt(8))

add_accent_line(slide, Inches(0.8), Inches(4.1), Inches(3.0), ACCENT_BLUE)

# Team Box
card = add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.2), fill_color=BG_CARD, border_color=BORDER_COLOR)
tf = add_text_box(slide, Inches(1.1), Inches(4.7), Inches(11.0), Inches(1.8))
set_para(tf, "PROJECT TEAM & EQUAL CONTRIBUTIONS", font_size=12, color=ACCENT_PURPLE, bold=True)

team_members = [
    ("Sagnik Basu", "23MID0042", "Architecture, Game Theory Engine & LLM Diagnostics"),
    ("C Sriharsha", "23MID0111", "Telemetry Ingestion Engine & Live Host/5G Collectors"),
    ("Maitree Singh", "23MID0076", "Generative Counterfactual Engine & Web App Dashboard"),
]

for i, (name, reg, role) in enumerate(team_members):
    x = Inches(1.1 + i * 3.7)
    tf_m = add_text_box(slide, x, Inches(5.2), Inches(3.5), Inches(1.2))
    set_para(tf_m, f"• {name} ({reg})", font_size=13, color=WHITE, bold=True)
    set_para(tf_m, role, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 2: PROJECT ABSTRACT & PROBLEM STATEMENT
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 2, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "OVERVIEW", bg_color=ACCENT_BLUE, width=Inches(1.4))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Abstract & Core Problem Statement", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_BLUE)

# Left Box: Abstract
card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=BORDER_COLOR)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "PROJECT ABSTRACT", font_size=14, color=ACCENT_BLUE, bold=True)
set_para(tf1, "Current predictive maintenance (PdM) operates as an opaque 'black box' — predicting failures without explaining why or how to fix them.", font_size=11, color=TEXT_SECONDARY, is_first=False, space_before=Pt(6))
set_para(tf1, "CTG-CPM is a closed-loop predictive maintenance recommendation system combining:", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf1, "1. Time-Series Diffusion Models for counterfactual 'What-If' future telemetry.", font_size=11, color=ACCENT_PURPLE, is_first=False)
set_para(tf1, "2. Graph Neural Networks (GraphSAGE) for topology-aware anomaly detection.", font_size=11, color=ACCENT_GREEN, is_first=False)
set_para(tf1, "3. Algorithmic Game Theory (VCG, Shapley, Nash) for provable agent coordination.", font_size=11, color=ACCENT_ORANGE, is_first=False)
set_para(tf1, "4. LLM-Powered Diagnostics (Groq GPT-OSS) for step-by-step remediation.", font_size=11, color=ACCENT_CYAN, is_first=False)

# Right Box: 3 Core Shortcomings
card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_RED)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "THREE CRITICAL INDUSTRY SHORTCOMINGS", font_size=14, color=ACCENT_RED, bold=True)

issues = [
    ("1. Opacity ('Black Box')", "Predicts failure timestamps but fails to explain the causal chain, leaving NOC engineers stranded."),
    ("2. Manual Bottleneck", "Engineers guess remediation scripts manually on live traffic, risking catastrophic cascading outages."),
    ("3. OPEX Truck Roll Drain", "Companies resort to costly physical hardware replacements because they cannot safely test software fixes."),
]
for title, desc in issues:
    set_para(tf2, title, font_size=12, color=WHITE, bold=True, is_first=False, space_before=Pt(10))
    set_para(tf2, desc, font_size=10, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 3: DETAILED SYSTEM ARCHITECTURE (5-LAYER STACK)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 3, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ARCHITECTURE", bg_color=ACCENT_PURPLE, width=Inches(1.6))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "5-Layer System Architecture & Data Flow", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_PURPLE)

layers = [
    ("LAYER 1: DATA INGESTION", "Apache Kafka, gNMI, psutil Live Telemetry & InfluxDB Storage", ACCENT_BLUE),
    ("LAYER 2: TOPOLOGY GNN", "GraphSAGE Topology Mapper & Anomaly Propagation Predictor", ACCENT_GREEN),
    ("LAYER 3: GenAI COUNTERFACTUAL", "Time-Series Diffusion Model generating parallel 'What-If' futures", ACCENT_PURPLE),
    ("LAYER 4: AGENTIC GAME THEORY", "Multi-Agent Loop: VCG Auction, Shapley RCA, Nash Equilibrium & SPE Solver", ACCENT_ORANGE),
    ("LAYER 5: REMEDIATION RECOMMENDATION", "Counterfactual projection selection; generates a recommended command (NOT auto-deployed)", ACCENT_CYAN),
]

for i, (lname, ldesc, color) in enumerate(layers):
    y = Inches(1.8 + i * 1.0)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), fill_color=BG_CARD, border_color=color)
    add_circle_number(slide, Inches(1.0), y + Inches(0.18), Inches(0.48), str(i+1), color)
    tf_l = add_text_box(slide, Inches(1.7), y + Inches(0.15), Inches(10.5), Inches(0.65))
    set_para(tf_l, lname, font_size=13, color=color, bold=True)
    set_para(tf_l, ldesc, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 4: DETAILED PROCESS PIPELINE & FLOWCHART
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 4, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "METHODOLOGY", bg_color=ACCENT_CYAN, width=Inches(1.6))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Detailed 4-Stage Process Pipeline", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_CYAN)

stages = [
    ("STAGE 1: DETECT", "GNN maps live topology and detects telemetry anomalies (OSNR drops, thermal spikes).", ACCENT_BLUE),
    ("STAGE 2: GENERATE", "Diffusion model generates N synthetic counterfactual future telemetry streams under candidate interventions.", ACCENT_PURPLE),
    ("STAGE 3: SELECT", "Agent 1 diagnoses root cause via Shapley Values; Agent 2 runs SPE Game & VCG auction; Agent 3 selects candidate via counterfactual projection.", ACCENT_GREEN),
    ("STAGE 4: RECOMMEND (NOT DEPLOYED)", "Produces a recommended NETCONF/YANG & PowerShell command. Auto-deployment is disabled in the prototype.", ACCENT_ORANGE),
]

for i, (stitle, sdesc, color) in enumerate(stages):
    x = Inches(0.8 + i * 2.95)
    card = add_shape(slide, x, Inches(2.0), Inches(2.75), Inches(4.8), fill_color=BG_CARD, border_color=color)
    add_pill_badge(slide, x + Inches(0.2), Inches(2.2), f"PHASE 0{i+1}", bg_color=color, text_color=BG_DARK, width=Inches(1.2))
    tf_s = add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.35), Inches(3.9))
    set_para(tf_s, stitle, font_size=13, color=WHITE, bold=True)
    set_para(tf_s, sdesc, font_size=10, color=TEXT_SECONDARY, is_first=False, space_before=Pt(8))

# ====================================================================
# SLIDE 5: DEEP DIVE — ALGORITHMS (PART 1: CORE ML STACK)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 5, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ALGORITHMS", bg_color=ACCENT_BLUE, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Algorithm Choice: Core ML & Generative Stack", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_BLUE)

# Box 1: Diffusion Model
card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "ALGORITHM 1: TIME-SERIES DIFFUSION MODEL", font_size=13, color=ACCENT_PURPLE, bold=True)
set_para(tf1, "Counterfactual Telemetry Generation (Diffusion-TS)", font_size=10, color=TEXT_DIM, is_first=False)
set_para(tf1, "• Method: Denoising score matching on multivariate time-series data.", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(8))
set_para(tf1, "• Intervention Conditioning: Synthesizes 'What-If' streams conditioned on specific actions (e.g. throttle CPU 15% vs increase cooling).", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "• Protocol Constraints: Ensures generated data obeys physical network laws (valid OSNR, laser bias, packet loss ratios).", font_size=11, color=TEXT_SECONDARY, is_first=False)

# Box 2: GraphSAGE GNN
card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "ALGORITHM 2: GRAPHSAGE GRAPH NEURAL NETWORK", font_size=13, color=ACCENT_GREEN, bold=True)
set_para(tf2, "Topology-Aware Anomaly Detection", font_size=10, color=TEXT_DIM, is_first=False)
set_para(tf2, "• Spatial-Temporal Learning: Learns node embeddings capturing both inter-device graph topology and time-series metric streams.", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(8))
set_para(tf2, "• Anomaly Propagation Modeling: Predicts downstream cascade failure risk when an optical node or CPU core degrades.", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf2, "• Dynamic Graphs: Topology refreshed periodically to keep node embeddings current.", font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 6: DEEP DIVE — ALGORITHMS (PART 2: ALGORITHMIC GAME THEORY)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 6, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "GAME THEORY", bg_color=ACCENT_ORANGE, text_color=BG_DARK, width=Inches(1.6))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Algorithmic Game Theory & Multi-Agent Stack", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_ORANGE)

gt_items = [
    ("1. VCG Auction Task Allocation", "Agents bid capability scores on tasks. VCG mechanism guarantees DSIC (Dominant-Strategy Incentive Compatibility) — truthfulness is optimal.", ACCENT_BLUE),
    ("2. Shapley Value Attribution", "Axiomatic feature attribution (Efficiency, Symmetry, Linearity, Null-Player) for fair multi-metric root-cause decomposition (e.g. CPU 50%, Temp 30%).", ACCENT_GREEN),
    ("3. SPE Backward Induction Game", "Solves extensive-form investment & bargaining game tree. Proves SPE strategy: (Low Investment, Greedy Proposal, Accept) -> Payoffs (13, 0).", ACCENT_PURPLE),
    ("4. Nash Equilibrium Coordination", "2-Player normal form matrix solver ensures multi-agent remediation strategies converge to a stable Nash Equilibrium, preventing conflicting actions.", ACCENT_ORANGE),
]

for i, (gtitle, gdesc, color) in enumerate(gt_items):
    y = Inches(1.8 + i * 1.3)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), fill_color=BG_CARD, border_color=color)
    tf_g = add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(11.1), Inches(0.95))
    set_para(tf_g, gtitle, font_size=13, color=color, bold=True)
    set_para(tf_g, gdesc, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 7: LLM-POWERED DIAGNOSTICS & EXPLANATION ENGINE
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 7, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "LLM ENGINE", bg_color=ACCENT_CYAN, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "LLM-Powered Intelligent Diagnostics (llm_diagnostician.py)", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_CYAN)

card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=BORDER_COLOR)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "GROQ API INTEGRATION (gpt-oss-20b)", font_size=13, color=ACCENT_CYAN, bold=True)
set_para(tf1, "• Role: Transforms low-level telemetry, Shapley weights, and game outcomes into structured JSON diagnostics.", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf1, "• Outputs Generated:", font_size=11, color=WHITE, bold=True, is_first=False, space_before=Pt(6))
set_para(tf1, "  - Human-Readable Problem Title & Description", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "  - Physical Root Cause Explanation", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "  - Step-by-Step Remediation Plan with impact", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "  - Risk Assessment if Unresolved within 1-4 hrs", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "  - Business Outcomes (Truck roll avoided, fix time)", font_size=11, color=TEXT_SECONDARY, is_first=False)

card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "SAMPLE LLM DIAGNOSTIC OUTPUT (JSON)", font_size=13, color=ACCENT_GREEN, bold=True)
code_sample = (
    "{\n"
    "  'severity': 'High',\n"
    "  'problem_title': 'CPU and Memory Saturation',\n"
    "  'root_cause': 'High CPU workload (50% Shapley) combined with memory pressure',\n"
    "  'risk_if_unresolved': 'Thermal throttling & system crash within 2-4 hours',\n"
    "  'remediation_steps': [\n"
    "    {'step': 1, 'action': 'Identify High CPU Procs'},\n"
    "    {'step': 2, 'action': 'Apply CPU Throttle 85%'},\n"
    "    {'step': 3, 'action': 'Optimize Memory Cache'}\n"
    "  ],\n"
    "  'truck_roll_avoided': 'potential (not measured)',\n"
    "  'estimated_fix_time': 'decision compute only (excludes deployment)'\n"
    "}"
)
set_para(tf2, code_sample, font_size=10, color=ACCENT_GREEN, is_first=False, space_before=Pt(6))

# ====================================================================
# SLIDE 8: ACTUAL WORKING PROGRESS SINCE R1
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 8, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "PROGRESS", bg_color=ACCENT_GREEN, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Actual Working Progress Since Review 1", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_GREEN)

progress_items = [
    ("Fully Built Core Engine", "Implemented all Python backend modules: game_engine.py, telemetry_collector.py, counterfactual_engine.py, agentic_remediator.py.", ACCENT_BLUE),
    ("Groq LLM Integration", "Created llm_diagnostician.py using openai/gpt-oss-20b model for structured, human-readable problem diagnosis and remediation.", ACCENT_PURPLE),
    ("HMW-Inspired Web Dashboard", "Developed modern glassmorphism web application (app.py + index.html) with real-time graphs, anomaly injection, and diagnostic panel.", ACCENT_GREEN),
    ("Automated Test Suite", "test_prototype.py unit/integration tests covering 5 core modules (Kafka test asserts honest status rather than emulating a broker).", ACCENT_ORANGE),
]

for i, (ptitle, pdesc, color) in enumerate(progress_items):
    y = Inches(1.8 + i * 1.3)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), fill_color=BG_CARD, border_color=color)
    tf_p = add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(11.1), Inches(0.95))
    set_para(tf_p, f"✓  {ptitle}", font_size=13, color=color, bold=True)
    set_para(tf_p, pdesc, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 9: CODE COMMIT & REPOSITORY EVIDENCE
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 9, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "EVIDENCE", bg_color=ACCENT_BLUE, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Code Committed & GitHub Repository Evidence", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_BLUE)

card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=BORDER_COLOR)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "PUBLIC GITHUB REPOSITORY", font_size=13, color=ACCENT_BLUE, bold=True)
set_para(tf1, "• Repository URL:", font_size=11, color=WHITE, bold=True, is_first=False, space_before=Pt(6))
set_para(tf1, "  github.com/sagnikbasutaan2004-source/\n  CTG-CPM-Self-Healing-Networks", font_size=10, color=ACCENT_CYAN, is_first=False)
set_para(tf1, "• Total Committed Files: 15 production files", font_size=11, color=TEXT_SECONDARY, is_first=False, space_before=Pt(6))
set_para(tf1, "• Total Lines of Code: 3,368+ lines", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "• Security Compliance: Zero API key leakage; GROQ_API_KEY managed via user environment variable.", font_size=11, color=ACCENT_GREEN, is_first=False, space_before=Pt(6))

card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "EXECUTABLE SYSTEM COMPONENTS", font_size=13, color=ACCENT_PURPLE, bold=True)

files = [
    ("game_engine.py", "Backward Induction SPE, VCG, Shapley, Nash"),
    ("telemetry_collector.py", "Live psutil host + 5G transceiver metrics"),
    ("counterfactual_engine.py", "'What-If' time-series future generator"),
    ("agentic_remediator.py", "3-Agent closed-loop orchestrator"),
    ("llm_diagnostician.py", "Groq LLM structured diagnosis generator"),
    ("app.py & templates/index.html", "Flask web app & glassmorphism UI"),
    ("main.py", "CLI orchestrator (in-process decision compute)"),
]
for fname, fdesc in files:
    set_para(tf2, f"• {fname}", font_size=10, color=WHITE, bold=True, is_first=False, space_before=Pt(4))
    set_para(tf2, f"   {fdesc}", font_size=9, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 10: JUSTIFICATION FOR DESIGN CHOICES
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 10, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "JUSTIFICATION", bg_color=ACCENT_PURPLE, width=Inches(1.8))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Justification for Design & Algorithmic Choices", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_PURPLE)

choices = [
    ("Time-Series Diffusion vs. GANs", "Diffusion models avoid mode collapse, provide superior sample quality, and support conditional generation on intervention variables.", ACCENT_BLUE),
    ("GNN vs. Standard CNNs", "Network topology is inherently non-Euclidean; our GNN (GCN or GraphSAGE, reported per run) models spatial node connectivity and cascading failure propagation.", ACCENT_GREEN),
    ("VCG & Nash Theory vs. Heuristic Prompting", "Game theory provides formal guarantees (DSIC truthfulness, stable equilibrium) that heuristic prompting lacks.", ACCENT_PURPLE),
    ("Groq Cloud LPU vs. Local LLMs", "Groq LPU offers fast LLM inference useful for the diagnostic step; it is additive to (not the whole) end-to-end MTTR.", ACCENT_ORANGE),
]

for i, (ctitle, cdesc, color) in enumerate(choices):
    y = Inches(1.8 + i * 1.3)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), fill_color=BG_CARD, border_color=color)
    tf_c = add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(11.1), Inches(0.95))
    set_para(tf_c, f"WHY {ctitle.upper()}?", font_size=13, color=color, bold=True)
    set_para(tf_c, cdesc, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 11: UPDATED FEASIBILITY & ECONOMIC PICTURE
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 11, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "FEASIBILITY", bg_color=ACCENT_ORANGE, text_color=BG_DARK, width=Inches(1.6))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Updated Economic Feasibility & Cost Analysis", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_ORANGE)

card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "POTENTIAL OPEX & HARDWARE COST", font_size=13, color=ACCENT_GREEN, bold=True)
set_para(tf1, "• Potential Truck-Roll Reduction: May avoid physical hardware replacement by recommending configuration fixes (potential benefit, not measured).", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf1, "• Possible Hardware Life Extension: E.g., optical laser bias adjustment might extend transceiver lifespan (not yet validated live).", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "• No New Sensor Cost: Uses existing telemetry streams (SNMP, gNMI, psutil) without new hardware sensors.", font_size=11, color=TEXT_SECONDARY, is_first=False)

card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_BLUE)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "CLOUD API & INFRASTRUCTURE COST", font_size=13, color=ACCENT_BLUE, bold=True)
set_para(tf2, "• Lightweight Inference: Distilled diffusion model and optimized GNN aim to keep GPU compute costs low (figure not yet benchmarked).", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf2, "• ROI: Cost-benefit depends on prevented outages; no dollar savings figure is claimed without live data.", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf2, "• Production Feasibility: 5-layer architecture uses proven open-source components (Kafka, PyTorch, Flask, psutil).", font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 12: SOCIAL & ENVIRONMENTAL CONSIDERATIONS
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 12, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "IMPACT", bg_color=ACCENT_GREEN, width=Inches(1.4))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Social & Environmental Considerations", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_GREEN)

impacts = [
    ("Potential Carbon Footprint Reduction", "Potentially fewer technician truck rolls could reduce fuel use from service vehicles (not yet measured).", ACCENT_GREEN),
    ("E-Waste Mitigation (potential)", "Software-based remediation could extend hardware life cycles, delaying premature e-waste generation.", ACCENT_CYAN),
    ("Reliability Support (potential)", "Fast remediation recommendations could support higher SLA uptime; end-to-end MTTR is not yet claimed.", ACCENT_BLUE),
    ("Human Operator Support", "Transitions NOC engineers from manual alert firefighting toward overseeing recommended remediations.", ACCENT_PURPLE),
]

for i, (ititle, idesc, color) in enumerate(impacts):
    y = Inches(1.8 + i * 1.3)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), fill_color=BG_CARD, border_color=color)
    tf_i = add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(11.1), Inches(0.95))
    set_para(tf_i, f"🌱  {ititle}", font_size=13, color=color, bold=True)
    set_para(tf_i, idesc, font_size=11, color=TEXT_SECONDARY, is_first=False)

# ====================================================================
# SLIDE 13: PROBLEMS ENCOUNTERED & TEAM ADAPTATION
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 13, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "ADAPTATION", bg_color=ACCENT_RED, width=Inches(1.6))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Problems Encountered & How Team Adapted", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_RED)

problems = [
    ("Problem 1: LLM Reasoning Syntax & Think Tags", "Qwen/Llama models initially included internal thinking tags (<think>...</think>) breaking raw JSON parsing.", "Adapted parser regex in llm_diagnostician.py & switched to openai/gpt-oss-20b for clean JSON output.", ACCENT_RED),
    ("Problem 2: VCG Computation Complexity", "Calculating full VCG auction social welfare over arbitrary agents can scale exponentially.", "Restricted auction search space to active 3-agent domain & implemented fast combinatorial allocation.", ACCENT_ORANGE),
    ("Problem 3: Simulation-Reality Gap", "Synthetic counterfactual futures might deviate from live host behavior.", "Added domain randomization and integrated real psutil telemetry baseline to constrain generation.", ACCENT_CYAN),
]

for i, (title, issue, adaptation, color) in enumerate(problems):
    y = Inches(1.8 + i * 1.7)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.55), fill_color=BG_CARD, border_color=color)
    tf_pr = add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(11.1), Inches(1.35))
    set_para(tf_pr, title, font_size=13, color=color, bold=True)
    set_para(tf_pr, f"• Issue: {issue}", font_size=10, color=TEXT_SECONDARY, is_first=False)
    set_para(tf_pr, f"• Team Adaptation: {adaptation}", font_size=10, color=ACCENT_GREEN, is_first=False)

# ====================================================================
# SLIDE 14: EQUAL INDIVIDUAL CONTRIBUTIONS DISTRIBUTION
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 14, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "TEAM ROLES", bg_color=ACCENT_PURPLE, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Equal Individual Contribution Breakdown", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_PURPLE)

members = [
    ("Sagnik Basu", "23MID0042", "33.3% Contribution",
     "• System Architecture & Patent Claim Formulation\n"
     "• Game Theory Engine (game_engine.py: VCG, Shapley, SPE, Nash)\n"
     "• Groq LLM Diagnostics Engine (llm_diagnostician.py)\n"
     "• Git Repository Management & Automated Testing", ACCENT_BLUE),
    ("C Sriharsha", "23MID0111", "33.3% Contribution",
     "• Telemetry Collector Engine (telemetry_collector.py)\n"
     "• Real-Time psutil Host Metric Ingestion (CPU, RAM, Disk, Processes)\n"
     "• 5G Optical Backhaul Transceiver Anomaly Simulator\n"
     "• CLI Controller & Main Pipeline Integration (main.py)", ACCENT_GREEN),
    ("Maitree Singh", "23MID0076", "33.3% Contribution",
     "• Generative Counterfactual Engine (counterfactual_engine.py)\n"
     "• Multi-Agent Remediation Orchestration (agentic_remediator.py)\n"
     "• HMW-Style Modern Glassmorphism Web App (templates/index.html)\n"
     "• Flask Server REST APIs & UI Visualizations (app.py)", ACCENT_ORANGE),
]

for i, (mname, mreg, mpct, mtasks, color) in enumerate(members):
    x = Inches(0.8 + i * 3.95)
    card = add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(5.0), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_pill_badge(slide, x + Inches(0.2), Inches(2.0), mpct, bg_color=color, text_color=BG_DARK, width=Inches(1.8))
    tf_m = add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(4.2))
    set_para(tf_m, mname, font_size=15, color=WHITE, bold=True)
    set_para(tf_m, f"Reg No: {mreg}", font_size=11, color=color, is_first=False)
    set_para(tf_m, mtasks, font_size=10, color=TEXT_SECONDARY, is_first=False, space_before=Pt(8))

# ====================================================================
# SLIDE 15: CONCLUSION & NEXT STEPS
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_footer_bar(slide, 15, TOTAL_SLIDES)

add_pill_badge(slide, Inches(0.8), Inches(0.5), "CONCLUSION", bg_color=ACCENT_GREEN, width=Inches(1.5))
tf = add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6))
set_para(tf, "Summary of Review 2 Achievements & Next Steps", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(2.0), ACCENT_GREEN)

card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)
tf1 = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
set_para(tf1, "REVIEW 2 ACHIEVEMENTS", font_size=13, color=ACCENT_GREEN, bold=True)
set_para(tf1, "✓ Functional Prototype with fast in-process decision compute (< a few ms, excludes LLM/deployment).", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf1, "✓ 15 Production Files committed to public GitHub repository.", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "✓ Automated Test Suite (5/5 unit & integration tests).", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf1, "✓ LLM-Powered Remediation Recommendation Engine via Groq API.", font_size=11, color=TEXT_SECONDARY, is_first=False)

card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), fill_color=BG_CARD, border_color=ACCENT_BLUE)
tf2 = add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
set_para(tf2, "NEXT STEPS (ROADMAP TO FINAL REVIEW)", font_size=13, color=ACCENT_BLUE, bold=True)
set_para(tf2, "1. Simu5G / ns-3 Sandbox Deployment: Scale testbed to 100+ optical transceiver nodes.", font_size=11, color=TEXT_PRIMARY, is_first=False, space_before=Pt(6))
set_para(tf2, "2. Provisional Patent Application (PPA): Complete filing for 3 novel patent claims.", font_size=11, color=TEXT_SECONDARY, is_first=False)
set_para(tf2, "3. Live Enterprise Pilot: Graduate system autonomy from advisory to semi-autonomous execution.", font_size=11, color=TEXT_SECONDARY, is_first=False)

# Save presentation
output_path = os.environ.get("OUTPUT_PPT_PATH", r"d:\Predictive Maintenance Project 3\CTG-CPM_Review2_Presentation.pptx")
prs.save(output_path)
print(f"Presentation generated successfully: {output_path}")
